"""
Enterprise AI Assistant - Document Processor Module
GeliÅŸmiÅŸ dÃ¶kÃ¼man iÅŸleme

Ã‡oklu format desteÄŸi ve akÄ±llÄ± Ã§Ä±karÄ±m.
"""

import re
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum


class DocumentType(str, Enum):
    """DÃ¶kÃ¼man tÃ¼rleri."""
    PDF = "pdf"
    WORD = "word"
    EXCEL = "excel"
    POWERPOINT = "powerpoint"
    TEXT = "text"
    MARKDOWN = "markdown"
    CODE = "code"
    HTML = "html"
    JSON = "json"
    CSV = "csv"
    IMAGE = "image"
    UNKNOWN = "unknown"


@dataclass
class DocumentMetadata:
    """DÃ¶kÃ¼man meta verileri."""
    filename: str
    file_path: str
    file_size: int
    file_hash: str
    doc_type: DocumentType
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    char_count: Optional[int] = None
    language: Optional[str] = None
    created_at: datetime = field(default_factory=datetime.now)
    processed_at: Optional[datetime] = None
    custom_metadata: Dict[str, Any] = field(default_factory=dict)
    
    def to_dict(self) -> Dict:
        return {
            "filename": self.filename,
            "file_path": self.file_path,
            "file_size": self.file_size,
            "file_hash": self.file_hash,
            "doc_type": self.doc_type.value,
            "page_count": self.page_count,
            "word_count": self.word_count,
            "char_count": self.char_count,
            "language": self.language,
            "created_at": self.created_at.isoformat(),
            "processed_at": self.processed_at.isoformat() if self.processed_at else None,
            "custom_metadata": self.custom_metadata,
        }


@dataclass
class ProcessedDocument:
    """Ä°ÅŸlenmiÅŸ dÃ¶kÃ¼man."""
    metadata: DocumentMetadata
    content: str
    sections: List[Dict[str, Any]] = field(default_factory=list)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)
    links: List[str] = field(default_factory=list)
    entities: List[Dict[str, str]] = field(default_factory=list)
    summary: Optional[str] = None


class DocumentProcessor:
    """GeliÅŸmiÅŸ dÃ¶kÃ¼man iÅŸleyici."""
    
    # Dosya uzantÄ±sÄ± -> DocumentType eÅŸleÅŸtirmesi
    EXTENSION_MAP = {
        ".pdf": DocumentType.PDF,
        ".doc": DocumentType.WORD,
        ".docx": DocumentType.WORD,
        ".xls": DocumentType.EXCEL,
        ".xlsx": DocumentType.EXCEL,
        ".ppt": DocumentType.POWERPOINT,
        ".pptx": DocumentType.POWERPOINT,
        ".txt": DocumentType.TEXT,
        ".md": DocumentType.MARKDOWN,
        ".markdown": DocumentType.MARKDOWN,
        ".py": DocumentType.CODE,
        ".js": DocumentType.CODE,
        ".ts": DocumentType.CODE,
        ".java": DocumentType.CODE,
        ".cpp": DocumentType.CODE,
        ".c": DocumentType.CODE,
        ".html": DocumentType.HTML,
        ".htm": DocumentType.HTML,
        ".json": DocumentType.JSON,
        ".csv": DocumentType.CSV,
        ".png": DocumentType.IMAGE,
        ".jpg": DocumentType.IMAGE,
        ".jpeg": DocumentType.IMAGE,
        ".gif": DocumentType.IMAGE,
    }
    
    def __init__(self):
        """Document processor baÅŸlat."""
        self._processed_cache: Dict[str, ProcessedDocument] = {}
    
    def detect_type(self, file_path: Path) -> DocumentType:
        """
        Dosya tÃ¼rÃ¼nÃ¼ tespit et.
        
        Args:
            file_path: Dosya yolu
            
        Returns:
            DÃ¶kÃ¼man tÃ¼rÃ¼
        """
        ext = file_path.suffix.lower()
        return self.EXTENSION_MAP.get(ext, DocumentType.UNKNOWN)
    
    def compute_hash(self, file_path: Path) -> str:
        """
        Dosya hash'i hesapla.
        
        Args:
            file_path: Dosya yolu
            
        Returns:
            MD5 hash
        """
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def extract_metadata(self, file_path: Path) -> DocumentMetadata:
        """
        Meta verileri Ã§Ä±kar.
        
        Args:
            file_path: Dosya yolu
            
        Returns:
            DÃ¶kÃ¼man meta verileri
        """
        path = Path(file_path)
        
        return DocumentMetadata(
            filename=path.name,
            file_path=str(path),
            file_size=path.stat().st_size,
            file_hash=self.compute_hash(path),
            doc_type=self.detect_type(path),
        )
    
    def extract_text_content(self, file_path: Path) -> str:
        """
        Metin iÃ§eriÄŸini Ã§Ä±kar.
        
        Args:
            file_path: Dosya yolu
            
        Returns:
            Metin iÃ§eriÄŸi
        """
        doc_type = self.detect_type(file_path)
        
        if doc_type == DocumentType.PDF:
            return self._extract_pdf(file_path)
        elif doc_type == DocumentType.WORD:
            return self._extract_word(file_path)
        elif doc_type == DocumentType.POWERPOINT:
            return self._extract_powerpoint(file_path)
        elif doc_type == DocumentType.EXCEL:
            return self._extract_excel(file_path)
        elif doc_type in [DocumentType.TEXT, DocumentType.MARKDOWN, DocumentType.CODE]:
            return self._extract_text(file_path)
        elif doc_type == DocumentType.HTML:
            return self._extract_html(file_path)
        elif doc_type == DocumentType.JSON:
            return self._extract_json(file_path)
        elif doc_type == DocumentType.CSV:
            return self._extract_csv(file_path)
        else:
            return self._extract_text(file_path)
    
    def _extract_pdf(self, file_path: Path) -> str:
        """PDF'den metin Ã§Ä±kar."""
        try:
            import pypdf
            
            text_parts = []
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text_parts.append(page.extract_text() or "")
            return "\n\n".join(text_parts)
        except ImportError:
            return "[PDF okuma iÃ§in pypdf gerekli]"
        except Exception as e:
            return f"[PDF okuma hatasÄ±: {e}]"
    
    def _extract_word(self, file_path: Path) -> str:
        """Word dÃ¶kÃ¼manÄ±ndan metin Ã§Ä±kar."""
        try:
            import docx
            
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs]
            return "\n\n".join(paragraphs)
        except ImportError:
            return "[Word okuma iÃ§in python-docx gerekli]"
        except Exception as e:
            return f"[Word okuma hatasÄ±: {e}]"
    
    def _extract_powerpoint(self, file_path: Path) -> str:
        """PowerPoint dÃ¶kÃ¼manÄ±ndan metin Ã§Ä±kar."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n## ðŸ“Š Slayt {slide_num}"]
                
                for shape in slide.shapes:
                    # BaÅŸlÄ±k
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text.append(text)
                    
                    # Tablolar
                    if shape.has_table:
                        table = shape.table
                        table_text = ["\n**Tablo:**"]
                        for row in table.rows:
                            row_cells = []
                            for cell in row.cells:
                                row_cells.append(cell.text.strip())
                            table_text.append(" | ".join(row_cells))
                        slide_text.extend(table_text)
                
                if len(slide_text) > 1:  # Sadece baÅŸlÄ±k deÄŸilse
                    text_parts.extend(slide_text)
            
            return "\n".join(text_parts)
        except ImportError:
            return "[PowerPoint okuma iÃ§in python-pptx gerekli]"
        except Exception as e:
            return f"[PowerPoint okuma hatasÄ±: {e}]"
    
    def _extract_excel(self, file_path: Path) -> str:
        """Excel'den metin Ã§Ä±kar - GeliÅŸmiÅŸ versiyon."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n## ðŸ“‘ Sayfa: {sheet_name}")
                
                # Ä°lk satÄ±rÄ± baÅŸlÄ±k olarak al
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    continue
                
                # BaÅŸlÄ±k satÄ±rÄ± varsa
                headers = rows[0] if rows else []
                has_header = headers and any(isinstance(h, str) for h in headers if h)
                
                if has_header:
                    header_text = " | ".join(str(h or "-") for h in headers)
                    text_parts.append(f"\n**BaÅŸlÄ±klar:** {header_text}")
                    text_parts.append("-" * 50)
                    data_rows = rows[1:]
                else:
                    data_rows = rows
                
                # Veri satÄ±rlarÄ±
                row_count = 0
                for row in data_rows:
                    row_text = " | ".join(str(cell or "-") for cell in row)
                    if row_text.strip(" |-"):
                        text_parts.append(row_text)
                        row_count += 1
                        # Ã‡ok fazla satÄ±r varsa Ã¶zetle
                        if row_count >= 100:
                            remaining = len(data_rows) - row_count
                            if remaining > 0:
                                text_parts.append(f"\n... ve {remaining} satÄ±r daha")
                            break
                
                text_parts.append(f"\nðŸ“Š Toplam: {len(data_rows)} satÄ±r veri")
            
            return "\n".join(text_parts)
        except ImportError:
            return "[Excel okuma iÃ§in openpyxl gerekli]"
        except Exception as e:
            return f"[Excel okuma hatasÄ±: {e}]"
    
    def _extract_text(self, file_path: Path) -> str:
        """DÃ¼z metin dosyasÄ± oku."""
        encodings = ["utf-8", "utf-16", "latin-1", "cp1254"]
        
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        
        return "[Dosya okunamadÄ±]"
    
    def _extract_html(self, file_path: Path) -> str:
        """HTML'den metin Ã§Ä±kar."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            
            # Remove script and style elements
            for element in soup(["script", "style"]):
                element.decompose()
            
            return soup.get_text(separator="\n", strip=True)
        except ImportError:
            # Fallback without BeautifulSoup
            text = self._extract_text(file_path)
            text = re.sub(r"<[^>]+>", "", text)
            return text
        except Exception as e:
            return f"[HTML okuma hatasÄ±: {e}]"
    
    def _extract_json(self, file_path: Path) -> str:
        """JSON'dan metin Ã§Ä±kar."""
        import json
        
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            return json.dumps(data, ensure_ascii=False, indent=2)
        except Exception as e:
            return f"[JSON okuma hatasÄ±: {e}]"
    
    def _extract_csv(self, file_path: Path) -> str:
        """CSV'den metin Ã§Ä±kar."""
        import csv
        
        try:
            text_parts = []
            with open(file_path, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                for row in reader:
                    text_parts.append(" | ".join(row))
            return "\n".join(text_parts)
        except Exception as e:
            return f"[CSV okuma hatasÄ±: {e}]"
    
    def extract_links(self, text: str) -> List[str]:
        """
        Metinden URL'leri Ã§Ä±kar.
        
        Args:
            text: Metin
            
        Returns:
            URL listesi
        """
        url_pattern = r'https?://[^\s<>"{}|\\^`\[\]]+'
        return list(set(re.findall(url_pattern, text)))
    
    def extract_emails(self, text: str) -> List[str]:
        """
        Metinden e-posta adreslerini Ã§Ä±kar.
        
        Args:
            text: Metin
            
        Returns:
            E-posta listesi
        """
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        return list(set(re.findall(email_pattern, text)))
    
    def extract_sections(self, text: str) -> List[Dict[str, Any]]:
        """
        Metinden bÃ¶lÃ¼mleri Ã§Ä±kar.
        
        Args:
            text: Metin
            
        Returns:
            BÃ¶lÃ¼m listesi
        """
        sections = []
        
        # Markdown baÅŸlÄ±klarÄ±
        header_pattern = r'^(#{1,6})\s+(.+)$'
        
        lines = text.split("\n")
        current_section = {"title": "GiriÅŸ", "content": [], "level": 0}
        
        for line in lines:
            match = re.match(header_pattern, line)
            if match:
                # Save current section
                if current_section["content"]:
                    current_section["content"] = "\n".join(current_section["content"])
                    sections.append(current_section)
                
                # Start new section
                level = len(match.group(1))
                title = match.group(2).strip()
                current_section = {"title": title, "content": [], "level": level}
            else:
                current_section["content"].append(line)
        
        # Save last section
        if current_section["content"]:
            current_section["content"] = "\n".join(current_section["content"])
            sections.append(current_section)
        
        return sections
    
    def detect_language(self, text: str) -> str:
        """
        Dili tespit et (basit heuristik).
        
        Args:
            text: Metin
            
        Returns:
            Dil kodu (tr, en, vb.)
        """
        # TÃ¼rkÃ§e karakterler
        turkish_chars = set("Ã§ÄŸÄ±Ã¶ÅŸÃ¼Ã‡ÄžÄ°Ã–ÅžÃœ")
        turkish_count = sum(1 for c in text if c in turkish_chars)
        
        # Basit heuristik
        if turkish_count > len(text) * 0.01:
            return "tr"
        
        return "en"
    
    def count_words(self, text: str) -> int:
        """Kelime sayÄ±sÄ±nÄ± hesapla."""
        return len(text.split())
    
    def process(self, file_path: Path) -> ProcessedDocument:
        """
        DÃ¶kÃ¼manÄ± tam olarak iÅŸle.
        
        Args:
            file_path: Dosya yolu
            
        Returns:
            Ä°ÅŸlenmiÅŸ dÃ¶kÃ¼man
        """
        path = Path(file_path)
        file_hash = self.compute_hash(path)
        
        # Check cache
        if file_hash in self._processed_cache:
            return self._processed_cache[file_hash]
        
        # Extract metadata
        metadata = self.extract_metadata(path)
        
        # Extract content
        content = self.extract_text_content(path)
        
        # Enrich metadata
        metadata.word_count = self.count_words(content)
        metadata.char_count = len(content)
        metadata.language = self.detect_language(content)
        metadata.processed_at = datetime.now()
        
        # Create processed document
        doc = ProcessedDocument(
            metadata=metadata,
            content=content,
            sections=self.extract_sections(content),
            links=self.extract_links(content),
            entities=[
                {"type": "email", "value": email}
                for email in self.extract_emails(content)
            ],
        )
        
        # Cache
        self._processed_cache[file_hash] = doc
        
        return doc
    
    def get_supported_extensions(self) -> List[str]:
        """Desteklenen dosya uzantÄ±larÄ±nÄ± dÃ¶ndÃ¼r."""
        return list(self.EXTENSION_MAP.keys())


# Singleton instance
document_processor = DocumentProcessor()
