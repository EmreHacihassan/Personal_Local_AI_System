"""
Enterprise AI Assistant - Async Document Processor
Radikal Performans Optimizasyonu - Çoklu Format Desteği

Özellikler:
- Async/await ile paralel işleme
- Robust timeout yönetimi
- Gelişmiş hata toleransı
- PPTX, XLSX, PDF için optimize edilmiş işleyiciler
- Büyük dosya desteği (chunk streaming)
"""

import asyncio
import concurrent.futures
import os
import traceback
from pathlib import Path
from typing import List, Dict, Any, Optional, Callable, Tuple
from datetime import datetime
import hashlib
import time
import threading
from dataclasses import dataclass, field
import queue


@dataclass
class ProcessingResult:
    """Dosya işleme sonucu."""
    filename: str
    success: bool
    chunks_created: int = 0
    error: Optional[str] = None
    processing_time: float = 0.0
    document_id: Optional[str] = None


@dataclass
class ProcessingProgress:
    """İşleme ilerlemesi."""
    total_files: int
    processed_files: int
    current_file: str = ""
    status: str = "pending"  # pending, processing, completed, error
    results: List[ProcessingResult] = field(default_factory=list)


class RobustDocumentLoader:
    """
    Robust document loader - Her format için özel optimizasyon.
    Takılma önleyici timeout mekanizması.
    """
    
    SUPPORTED_EXTENSIONS = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "doc",
        ".pptx": "pptx",
        ".ppt": "ppt",
        ".xlsx": "xlsx",
        ".xls": "xls",
        ".csv": "csv",
        ".txt": "text",
        ".md": "markdown",
        ".html": "html",
        ".htm": "html",
        ".json": "json",
        ".xml": "xml",
    }
    
    # Format bazlı timeout (saniye)
    TIMEOUTS = {
        "pdf": 120,      # PDF'ler uzun sürebilir
        "pptx": 90,      # PowerPoint
        "xlsx": 90,      # Excel
        "docx": 60,      # Word
        "default": 60,
    }
    
    def __init__(self):
        self._executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
    
    def get_timeout(self, file_type: str, file_size_mb: float) -> int:
        """Dosya boyutu ve tipine göre dinamik timeout."""
        base_timeout = self.TIMEOUTS.get(file_type, self.TIMEOUTS["default"])
        # Her 10MB için 30 saniye ekle
        size_factor = int(file_size_mb / 10) * 30
        return base_timeout + size_factor
    
    def load_with_timeout(self, file_path: str, timeout: int = 60) -> Tuple[List[Dict], Optional[str]]:
        """
        Timeout ile dosya yükle.
        Returns: (documents, error)
        """
        path = Path(file_path)
        if not path.exists():
            return [], f"Dosya bulunamadı: {file_path}"
        
        extension = path.suffix.lower()
        if extension not in self.SUPPORTED_EXTENSIONS:
            return [], f"Desteklenmeyen format: {extension}"
        
        file_type = self.SUPPORTED_EXTENSIONS[extension]
        file_size_mb = path.stat().st_size / (1024 * 1024)
        
        # Dinamik timeout
        actual_timeout = self.get_timeout(file_type, file_size_mb)
        if timeout > 0:
            actual_timeout = min(actual_timeout, timeout)
        
        try:
            # Thread pool ile timeout korumalı işleme
            future = self._executor.submit(self._load_file_internal, path, file_type)
            documents = future.result(timeout=actual_timeout)
            return documents, None
        except concurrent.futures.TimeoutError:
            return [], f"Timeout ({actual_timeout}s): Dosya çok büyük veya karmaşık"
        except Exception as e:
            return [], f"Hata: {str(e)[:200]}"
    
    def _load_file_internal(self, path: Path, file_type: str) -> List[Dict]:
        """İç dosya yükleme - format bazlı."""
        metadata = self._get_base_metadata(path)
        
        try:
            if file_type == "pdf":
                return self._load_pdf_robust(path, metadata)
            elif file_type == "docx":
                return self._load_docx_robust(path, metadata)
            elif file_type in ("pptx", "ppt"):
                return self._load_pptx_robust(path, metadata)
            elif file_type in ("xlsx", "xls"):
                return self._load_excel_robust(path, metadata)
            elif file_type == "csv":
                return self._load_csv_robust(path, metadata)
            elif file_type in ("text", "markdown"):
                return self._load_text(path, metadata)
            elif file_type == "html":
                return self._load_html_robust(path, metadata)
            elif file_type == "json":
                return self._load_json(path, metadata)
            else:
                return self._load_text(path, metadata)
        except Exception as e:
            # Fallback: En azından dosya bilgisini döndür
            return [{
                "content": f"[İçerik Okunamadı]\n\nDosya: {path.name}\nTip: {file_type}\nHata: {str(e)[:200]}",
                "metadata": metadata
            }]
    
    def _get_base_metadata(self, path: Path) -> Dict[str, Any]:
        """Temel metadata."""
        stat = path.stat()
        return {
            "source": str(path),
            "filename": path.name,
            "file_type": path.suffix.lower(),
            "file_size": stat.st_size,
            "created_at": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        }
    
    def _load_pdf_robust(self, path: Path, metadata: Dict) -> List[Dict]:
        """PDF yükle - çok katmanlı fallback."""
        all_content = []
        
        # Yöntem 1: pypdf
        try:
            from pypdf import PdfReader
            import warnings
            with warnings.catch_warnings():
                warnings.filterwarnings("ignore")
                reader = PdfReader(str(path), strict=False)
                for i, page in enumerate(reader.pages, 1):
                    try:
                        text = page.extract_text() or ""
                        if text.strip():
                            all_content.append(f"[Sayfa {i}]\n{text.strip()}")
                    except:
                        continue
                
                if all_content:
                    metadata["total_pages"] = len(reader.pages)
                    metadata["method"] = "pypdf"
                    return [{"content": "\n\n".join(all_content), "metadata": metadata}]
        except:
            pass
        
        # Yöntem 2: pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    try:
                        text = page.extract_text() or ""
                        if text.strip():
                            all_content.append(f"[Sayfa {i}]\n{text.strip()}")
                        
                        # Tabloları da çıkar
                        tables = page.extract_tables()
                        for table in tables:
                            if table:
                                table_text = []
                                for row in table:
                                    if row:
                                        row_text = [str(cell) if cell else "" for cell in row]
                                        table_text.append(" | ".join(row_text))
                                if table_text:
                                    all_content.append(f"[Tablo - Sayfa {i}]\n" + "\n".join(table_text))
                    except:
                        continue
                
                if all_content:
                    metadata["total_pages"] = len(pdf.pages)
                    metadata["method"] = "pdfplumber"
                    return [{"content": "\n\n".join(all_content), "metadata": metadata}]
        except:
            pass
        
        # Fallback
        return [{
            "content": f"[PDF İçeriği Çıkarılamadı]\n\nDosya: {path.name}\nBoyut: {metadata.get('file_size', 0) / 1024:.1f} KB",
            "metadata": metadata
        }]
    
    def _load_docx_robust(self, path: Path, metadata: Dict) -> List[Dict]:
        """Word yükle - robust versiyon."""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(str(path))
            parts = []
            
            # Paragraflar
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text.strip())
            
            # Tablolar
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        table_text.append(" | ".join(cells))
                if table_text:
                    parts.append("\n[Tablo]\n" + "\n".join(table_text))
            
            content = "\n\n".join(parts) if parts else f"[Boş Döküman: {path.name}]"
            return [{"content": content, "metadata": metadata}]
            
        except Exception as e:
            return [{
                "content": f"[Word İçeriği Okunamadı]\n\nDosya: {path.name}\nHata: {str(e)[:100]}",
                "metadata": metadata
            }]
    
    def _load_pptx_robust(self, path: Path, metadata: Dict) -> List[Dict]:
        """PowerPoint yükle - gelişmiş versiyon."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            prs = Presentation(str(path))
            all_slides = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"\n## Slayt {slide_num}"]
                
                for shape in slide.shapes:
                    # Metin kutuları
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_parts.append(text)
                    
                    # Tablolar
                    if shape.has_table:
                        table_text = ["\n**Tablo:**"]
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            table_text.append(" | ".join(cells))
                        slide_parts.extend(table_text)
                    
                    # SmartArt ve diğer şekiller
                    try:
                        if hasattr(shape, 'text') and shape.text:
                            if shape.text.strip() not in [p for p in slide_parts]:
                                slide_parts.append(shape.text.strip())
                    except:
                        pass
                
                if len(slide_parts) > 1:
                    all_slides.extend(slide_parts)
            
            content = "\n".join(all_slides) if all_slides else f"[Boş Sunum: {path.name}]"
            metadata["total_slides"] = len(prs.slides)
            return [{"content": content, "metadata": metadata}]
            
        except Exception as e:
            return [{
                "content": f"[PowerPoint İçeriği Okunamadı]\n\nDosya: {path.name}\nHata: {str(e)[:100]}",
                "metadata": metadata
            }]
    
    def _load_excel_robust(self, path: Path, metadata: Dict) -> List[Dict]:
        """Excel yükle - çoklu sayfa desteği."""
        documents = []
        
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(str(path), data_only=True, read_only=True)
            
            for sheet_name in wb.sheetnames:
                try:
                    sheet = wb[sheet_name]
                    rows = []
                    row_count = 0
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_count += 1
                        if row_count > 10000:  # Çok büyük sayfaları sınırla
                            rows.append("[... Devamı kesildi (10000+ satır) ...]")
                            break
                        
                        cells = [str(cell) if cell is not None else "" for cell in row]
                        if any(c.strip() for c in cells):
                            rows.append(" | ".join(cells))
                    
                    if rows:
                        sheet_meta = metadata.copy()
                        sheet_meta["sheet_name"] = sheet_name
                        sheet_meta["row_count"] = row_count
                        
                        content = f"## Sayfa: {sheet_name}\n\n" + "\n".join(rows)
                        documents.append({"content": content, "metadata": sheet_meta})
                except Exception as e:
                    continue
            
            wb.close()
            
            if not documents:
                return [{
                    "content": f"[Boş Excel: {path.name}]",
                    "metadata": metadata
                }]
            
            return documents
            
        except Exception as e:
            return [{
                "content": f"[Excel İçeriği Okunamadı]\n\nDosya: {path.name}\nHata: {str(e)[:100]}",
                "metadata": metadata
            }]
    
    def _load_csv_robust(self, path: Path, metadata: Dict) -> List[Dict]:
        """CSV yükle."""
        try:
            import csv
            
            rows = []
            row_count = 0
            
            # Encoding detect
            encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-9']
            
            for encoding in encodings:
                try:
                    with open(path, "r", encoding=encoding, newline='') as f:
                        # Delimiter detect
                        sample = f.read(4096)
                        f.seek(0)
                        
                        delimiter = ','
                        if ';' in sample and sample.count(';') > sample.count(','):
                            delimiter = ';'
                        elif '\t' in sample and sample.count('\t') > sample.count(','):
                            delimiter = '\t'
                        
                        reader = csv.reader(f, delimiter=delimiter)
                        
                        for row in reader:
                            row_count += 1
                            if row_count > 10000:
                                rows.append("[... Devamı kesildi (10000+ satır) ...]")
                                break
                            if any(cell.strip() for cell in row):
                                rows.append(" | ".join(row))
                        break
                except UnicodeDecodeError:
                    continue
            
            content = "\n".join(rows) if rows else f"[Boş CSV: {path.name}]"
            metadata["row_count"] = row_count
            return [{"content": content, "metadata": metadata}]
            
        except Exception as e:
            return [{
                "content": f"[CSV İçeriği Okunamadı]\n\nDosya: {path.name}\nHata: {str(e)[:100]}",
                "metadata": metadata
            }]
    
    def _load_text(self, path: Path, metadata: Dict) -> List[Dict]:
        """Text/Markdown yükle."""
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-9']
        
        for encoding in encodings:
            try:
                with open(path, "r", encoding=encoding) as f:
                    content = f.read()
                return [{"content": content, "metadata": metadata}]
            except UnicodeDecodeError:
                continue
        
        # Binary fallback
        with open(path, "rb") as f:
            content = f.read().decode('utf-8', errors='ignore')
        return [{"content": content, "metadata": metadata}]
    
    def _load_html_robust(self, path: Path, metadata: Dict) -> List[Dict]:
        """HTML yükle."""
        try:
            from bs4 import BeautifulSoup
            
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            
            # Script ve style kaldır
            for elem in soup(["script", "style", "meta", "link"]):
                elem.decompose()
            
            content = soup.get_text(separator="\n", strip=True)
            
            title = soup.find("title")
            if title:
                metadata["title"] = title.get_text(strip=True)
            
            return [{"content": content, "metadata": metadata}]
            
        except Exception as e:
            return self._load_text(path, metadata)
    
    def _load_json(self, path: Path, metadata: Dict) -> List[Dict]:
        """JSON yükle."""
        try:
            import json
            
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
            
            content = json.dumps(data, indent=2, ensure_ascii=False)
            return [{"content": content, "metadata": metadata}]
            
        except Exception as e:
            return [{
                "content": f"[JSON İçeriği Okunamadı]\n\nDosya: {path.name}\nHata: {str(e)[:100]}",
                "metadata": metadata
            }]


class AsyncBatchProcessor:
    """
    Async batch document processor.
    Birden fazla dosyayı paralel işler.
    """
    
    def __init__(self, max_workers: int = 3):
        self.max_workers = max_workers
        self.loader = RobustDocumentLoader()
        self._progress: Dict[str, ProcessingProgress] = {}
        self._lock = threading.Lock()
    
    def get_progress(self, batch_id: str) -> Optional[ProcessingProgress]:
        """İşleme ilerlemesini al."""
        with self._lock:
            return self._progress.get(batch_id)
    
    async def process_files_async(
        self,
        file_paths: List[str],
        batch_id: str,
        on_file_complete: Optional[Callable[[ProcessingResult], None]] = None
    ) -> List[ProcessingResult]:
        """
        Dosyaları async olarak işle.
        """
        # Progress başlat
        with self._lock:
            self._progress[batch_id] = ProcessingProgress(
                total_files=len(file_paths),
                processed_files=0,
                status="processing"
            )
        
        results = []
        
        # Semaphore ile paralel işleme sınırla
        semaphore = asyncio.Semaphore(self.max_workers)
        
        async def process_single(file_path: str) -> ProcessingResult:
            async with semaphore:
                return await self._process_file(file_path, batch_id)
        
        # Tüm dosyaları paralel işle
        tasks = [process_single(fp) for fp in file_paths]
        
        for coro in asyncio.as_completed(tasks):
            result = await coro
            results.append(result)
            
            # Progress güncelle
            with self._lock:
                if batch_id in self._progress:
                    self._progress[batch_id].processed_files += 1
                    self._progress[batch_id].results.append(result)
            
            if on_file_complete:
                on_file_complete(result)
        
        # Tamamlandı
        with self._lock:
            if batch_id in self._progress:
                self._progress[batch_id].status = "completed"
        
        return results
    
    async def _process_file(self, file_path: str, batch_id: str) -> ProcessingResult:
        """Tek dosya işle."""
        start_time = time.time()
        filename = Path(file_path).name
        
        # Progress güncelle
        with self._lock:
            if batch_id in self._progress:
                self._progress[batch_id].current_file = filename
        
        try:
            # Thread pool'da dosya yükle (blocking I/O)
            loop = asyncio.get_event_loop()
            documents, error = await loop.run_in_executor(
                None,
                lambda: self.loader.load_with_timeout(file_path, timeout=180)
            )
            
            if error:
                return ProcessingResult(
                    filename=filename,
                    success=False,
                    error=error,
                    processing_time=time.time() - start_time
                )
            
            return ProcessingResult(
                filename=filename,
                success=True,
                chunks_created=len(documents),
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            return ProcessingResult(
                filename=filename,
                success=False,
                error=str(e)[:200],
                processing_time=time.time() - start_time
            )
    
    def cleanup(self, batch_id: str):
        """Batch progress temizle."""
        with self._lock:
            if batch_id in self._progress:
                del self._progress[batch_id]


# Singleton instances
robust_loader = RobustDocumentLoader()
batch_processor = AsyncBatchProcessor(max_workers=3)
